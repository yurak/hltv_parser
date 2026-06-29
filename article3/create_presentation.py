"""
Скрипт для створення наукової презентації з аналізу інваріантності ознак CS2/CSGO
Вимоги: pip install python-pptx Pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Кольорова схема
DARK_BLUE = RGBColor(0, 51, 102)
ACCENT_BLUE = RGBColor(0, 112, 192)
ACCENT_GREEN = RGBColor(0, 176, 80)
ACCENT_RED = RGBColor(255, 0, 0)
LIGHT_GRAY = RGBColor(242, 242, 242)
WHITE = RGBColor(255, 255, 255)


def add_title_slide(prs, title, subtitle=""):
    """Додає титульний слайд"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Фоновий прямокутник
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(10), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Заголовок
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points, image_path=None):
    """Додає слайд з контентом"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Заголовок слайду
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Лінія під заголовком
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(9), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Контент
    if image_path and os.path.exists(image_path):
        # Зображення справа, текст зліва
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(5))
        try:
            slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.4), width=Inches(4.5))
        except Exception as e:
            print(f"Не вдалося додати зображення {image_path}: {e}")
    else:
        # Тільки текст на всю ширину
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        p.level = 0

    return slide


def add_table_slide(prs, title, headers, rows):
    """Додає слайд з таблицею"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Таблиця
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5 * num_rows)
    ).table

    # Заголовки
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(14)

    # Дані
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)

    return slide


def add_results_slide(prs, title, percentage, invariant_count, total_count, top_invariant, top_variant, image_path=None):
    """Додає слайд з результатами"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Великий відсоток
    pct_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(1))
    tf = pct_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{percentage}%"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    # Підпис
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(3), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{invariant_count} з {total_count} ознак інваріантні"
    p.font.size = Pt(14)

    # Топ інваріантні
    inv_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(4.3), Inches(2))
    tf = inv_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Найстабільніші:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    for item in top_invariant:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(12)

    # Топ варіантні
    var_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4.3), Inches(2))
    tf = var_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Найваріантніші:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED

    for item in top_variant:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(12)

    # Зображення
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(5), Inches(1.2), width=Inches(4.7))
        except Exception as e:
            print(f"Не вдалося додати зображення {image_path}: {e}")

    return slide


def add_key_findings_slide(prs, title, features):
    """Слайд з ключовими ознаками"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Підзаголовок
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Стабільні по картах, версіях ТА сторонах"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ACCENT_BLUE

    # Ознаки у дві колонки
    left_features = features[:6]
    right_features = features[6:]

    # Ліва колонка
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    for i, (feature, desc) in enumerate(left_features):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"✓ {feature}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN

        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(12)
        p.space_after = Pt(8)

    # Права колонка
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    for i, (feature, desc) in enumerate(right_features):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = f"✓ {feature}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN

        p = tf.add_paragraph()
        p.text = f"   {desc}"
        p.font.size = Pt(12)
        p.space_after = Pt(8)

    return slide


def add_conclusion_slide(prs, conclusions):
    """Слайд висновків"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Висновки"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE

    # Висновки
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame

    for i, (num, text) in enumerate(conclusions):
        if i > 0:
            p = tf.add_paragraph()
            p.space_before = Pt(16)
        else:
            p = tf.paragraphs[0]

        # Номер
        run = p.add_run()
        run.text = f"{num}. "
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = ACCENT_BLUE

        # Текст
        run = p.add_run()
        run.text = text
        run.font.size = Pt(18)

    return slide


def create_presentation():
    """Головна функція створення презентації"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    base_path = os.path.dirname(os.path.abspath(__file__))

    # 1. Титульний слайд
    add_title_slide(
        prs,
        "Аналіз інваріантності ознак\nпродуктивності гравців CS2/CSGO",
        "Визначення контекстно-незалежних метрик скілу"
    )

    # 2. Актуальність
    add_content_slide(prs, "Актуальність проблеми", [
        "Кіберспорт — індустрія з оборотом понад $1.8 млрд",
        "Counter-Strike — одна з найпопулярніших кіберспортивних дисциплін",
        "Об'єктивна оцінка гравців критична для скаутингу та формування складів",
        "Традиційні метрики (K/D, ADR) залежать від контексту гри",
        "Проблема: як відокремити справжній скіл від ситуативних факторів?"
    ])

    # 3. Дослідницьке питання
    add_content_slide(prs, "Дослідницьке питання", [
        "Які метрики продуктивності є ІНВАРІАНТНИМИ?",
        "",
        "Три виміри аналізу:",
        "   1. Карти (8 різних карт)",
        "   2. Версії гри (CS2 vs CSGO)",
        "   3. Сторони (Терористи vs Контр-терористи)",
        "",
        "Інваріантні ознаки = справжній скіл гравця"
    ])

    # 4. Набір даних
    add_table_slide(prs, "Набір даних",
        ["Параметр", "Значення"],
        [
            ["Загальна кількість записів", "3 486"],
            ["Записів CS2", "1 820"],
            ["Записів CSGO", "1 666"],
            ["Кількість карт", "8"],
            ["Числових ознак", "104 → 36 (після відбору)"],
            ["Джерело", "HLTV.org"]
        ]
    )

    # 5. Методологія - загальний підхід
    add_content_slide(prs, "Методологія: чому розміри ефектів?", [
        "Проблема p-values при великих вибірках (n > 3000):",
        "   • Навіть мізерні відмінності 'статистично значущі'",
        "   • Не показують практичну важливість",
        "",
        "Рішення — розміри ефектів:",
        "   • η² (eta-squared) — частка поясненої дисперсії",
        "   • Cohen's d — стандартизована різниця середніх",
        "   • Масштабно-незалежні міри практичної значущості"
    ])

    # 6. Методологія - таблиця тестів
    add_table_slide(prs, "Методологія: статистичні тести",
        ["Аналіз", "Тест", "Розмір ефекту", "Поріг"],
        [
            ["Карти (8 груп)", "ANOVA", "η² (eta-squared)", "< 0.06"],
            ["Версії (2 групи)", "Mann-Whitney U", "Cohen's d", "< 0.5"],
            ["Сторони (парні)", "Парний t-тест", "Cohen's d", "< 0.5"]
        ]
    )

    # 7. Результати - карти
    add_results_slide(prs,
        "Результати: інваріантність по картах",
        75, 27, 36,
        ["opening_deaths_per_round (η² = 0.003)",
         "opening_attempts (η² = 0.005)",
         "opening_kills_per_round (η² = 0.008)"],
        ["utility_damage_per_round (η² = 0.156)",
         "flashes_thrown_per_round (η² = 0.142)",
         "utility (η² = 0.128)"],
        os.path.join(base_path, "normalized_map_heatmap.png")
    )

    # 8. Результати - версії
    add_results_slide(prs,
        "Результати: інваріантність по версіях",
        92, 33, 36,
        ["opening_deaths_per_round (d = 0.02)",
         "utility (d = 0.04)",
         "clutch_points_per_round (d = 0.06)"],
        ["assisted_kills_percentage (d = 0.89)",
         "assists_per_round (d = 0.78)",
         "support_rounds (d = 0.71)"],
        os.path.join(base_path, "normalized_version_comparison.png")
    )

    # 9. Результати - сторони
    add_results_slide(prs,
        "Результати: інваріантність по сторонах",
        48, 14, 29,
        ["opening_attempts (d = 0.03)",
         "entrying (d = 0.05)",
         "pistol_round_rating (d = 0.08)"],
        ["saves_per_round_loss (d = 1.24)",
         "utility_damage_per_round (d = 0.98)",
         "flashes_thrown_per_round (d = 0.94)"],
        os.path.join(base_path, "normalized_side_comparison.png")
    )

    # 10. Універсальні ознаки
    add_key_findings_slide(prs, "11 універсально інваріантних ознак", [
        ("rating_3.0", "Загальний рейтинг гравця"),
        ("firepower", "Здатність до фрагів"),
        ("opening", "Участь у відкриваючих дуелях"),
        ("opening_attempts", "Частота перших боїв"),
        ("pistol_round_rating", "Продуктивність у пістолетних раундах"),
        ("clutching", "Гра в клатч-ситуаціях"),
        ("clutch_points_per_round", "Внесок у клатчі"),
        ("entrying", "Схильність до entry-фрагів"),
        ("sniping", "Стиль гри з AWP"),
        ("sniper_multi_kill_rounds", "Раунди з мультікілами AWP"),
        ("sniper_opening_kills_per_round", "Опенінг кіли з AWP")
    ])

    # 11. PCA візуалізація
    add_content_slide(prs, "PCA візуалізація", [
        "Проекція 36-вимірного простору ознак у 2D",
        "",
        "Спостереження:",
        "   • Карти — значне перекриття, немає кластерів",
        "   • Версії — CS2 та CSGO майже повністю перекриваються",
        "   • Сторони — помітне розділення по PC2"
    ], os.path.join(base_path, "pca_visualization.png"))

    # 12. Практичні імплікації
    add_content_slide(prs, "Практичні імплікації", [
        "Для скаутингу талантів:",
        "   • Фокус на 11 універсальних ознаках",
        "   • Opening stats як індикатор механічного скілу",
        "",
        "Для ML-моделей:",
        "   • Інваріантні ознаки як предиктори",
        "   • Уникати utility-метрик для крос-контекстних моделей",
        "",
        "Для аналітики:",
        "   • Нормалізувати по стороні при порівнянні",
        "   • Враховувати map pool команди"
    ])

    # 13. Обмеження
    add_content_slide(prs, "Обмеження дослідження", [
        "Вибірка: тільки топ-гравці за рейтингом HLTV",
        "   → Результати можуть не узагальнюватись на аматорський рівень",
        "",
        "Агрегація: статистика по матчах",
        "   → Втрачається раунд-за-раундом динаміка",
        "",
        "Cross-sectional дизайн:",
        "   → Не аналізуємо темпоральну стабільність",
        "",
        "Не враховано ігрові ролі:",
        "   → AWPer vs Rifler vs IGL можуть мати різні патерни"
    ])

    # 14. Майбутні дослідження
    add_content_slide(prs, "Напрями подальших досліджень", [
        "Лонгітюдний аналіз",
        "   • Стабільність ознак протягом кар'єри гравця",
        "",
        "Аналіз по ролях",
        "   • Інваріантність для AWPers vs Riflers vs IGLs",
        "",
        "Багаторівневе моделювання",
        "   • Гравець → Команда → Турнір",
        "",
        "Предиктивна валідація",
        "   • Чи інваріантні ознаки прогнозують майбутній перформанс?"
    ])

    # 15. Висновки
    add_conclusion_slide(prs, [
        (1, "75% ознак інваріантні до карт — карта менше впливає на скіл"),
        (2, "92% ознак інваріантні до версії — базові навички переносяться"),
        (3, "48% ознак інваріантні до сторони — T/CT має найбільший вплив"),
        (4, "11 універсальних ознак відображають контекстно-незалежний скіл"),
        (5, "Utility-метрики найваріантніші — залежать від тактики")
    ])

    # 16. Фінальний слайд
    add_title_slide(prs, "Дякую за увагу!", "Питання?")

    # Зберігаємо
    output_path = os.path.join(base_path, "presentation_invariance_analysis.pptx")
    prs.save(output_path)
    print(f"Презентацію збережено: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
